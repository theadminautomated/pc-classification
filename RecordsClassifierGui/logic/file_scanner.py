#!/usr/bin/env python3
"""
File Scanner Module for Electronic Records Classification
---------------------------------------------------------
Efficiently discovers and categorizes files for processing with proper architecture.

Run as script to see file stats for a directory:
    python file_scanner.py /path/to/folder
"""

import os
import sys
import subprocess
from pathlib import Path
from typing import Dict, Set, List, Iterator, Tuple, Union
from dataclasses import dataclass
import datetime
import logging

try:
    from docx import Document
except Exception:
    Document = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import xlrd
except Exception:
    xlrd = None

try:
    from PIL import Image
    import pytesseract
except Exception:
    Image = None
    pytesseract = None

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("file_scanner")

INCLUDE_EXT: Set[str] = frozenset({
    '.txt', '.csv', '.docx', '.xlsx', '.pptx', '.pdf', '.html', '.htm', '.md',
    '.rtf', '.odt', '.xml', '.json', '.yaml', '.yml', '.log', '.tsv'
})

EXCLUDE_EXT: Set[str] = frozenset({
    '.tmp', '.bak', '.old', '.zip', '.rar', '.tar', '.gz', '.7z',
    '.exe', '.dll', '.sys', '.iso', '.dmg', '.apk', '.msi', '.ps1', '.psd1',
    '.psm1', '.db', '.mdb', '.accdb'
})

@dataclass
class FileInfo:
    """Information about a discovered file."""
    path: Path
    size_bytes: int
    modified_time: datetime.datetime
    extension: str
    category: str  # 'destroy', 'analyze', 'skip'
    reason: str

class FileScanner:
    """
    Handles file discovery and categorization for classification.
    """
    def __init__(self, include_ext: Set[str] = None, exclude_ext: Set[str] = None):
        self.include_ext = include_ext or INCLUDE_EXT
        self.exclude_ext = exclude_ext or EXCLUDE_EXT
        self.destroy_threshold = datetime.datetime.now() - datetime.timedelta(days=6 * 365)

    def scan_directory(self, directory_path: Union[str, Path]) -> Iterator[FileInfo]:
        """
        Scan a directory and yield FileInfo objects for all discovered files.
        """
        directory = Path(directory_path)
        if not directory.exists():
            raise ValueError(f"Directory does not exist: {directory}")
        if not directory.is_dir():
            raise ValueError(f"Path is not a directory: {directory}")
        logger.info(f"Scanning directory: {directory}")

        for file_path in directory.rglob('*'):
            if not file_path.is_file():
                continue
            if file_path.name.startswith('.') or file_path.name.startswith('~$'):
                continue
            try:
                file_info = self._analyze_file(file_path)
                yield file_info
            except Exception as e:
                logger.warning(f"Error analyzing file {file_path}: {e}")
                yield FileInfo(
                    path=file_path,
                    size_bytes=0,
                    modified_time=datetime.datetime.now(),
                    extension=file_path.suffix.lower(),
                    category='skip',
                    reason=f"Error analyzing file: {e}"
                )

    def _analyze_file(self, file_path: Path) -> FileInfo:
        """Analyze a single file and determine its category."""
        stat_info = file_path.stat()
        modified_time = datetime.datetime.fromtimestamp(stat_info.st_mtime)
        extension = file_path.suffix.lower()
        category, reason = self._categorize_file(modified_time, extension)
        return FileInfo(
            path=file_path,
            size_bytes=stat_info.st_size,
            modified_time=modified_time,
            extension=extension,
            category=category,
            reason=reason
        )

    def _categorize_file(self, modified_time: datetime.datetime, extension: str) -> Tuple[str, str]:
        """Categorize a file based on age and type."""
        if modified_time < self.destroy_threshold:
            return 'destroy', 'Older than 6 years - automatic destroy'
        if extension in self.exclude_ext:
            return 'skip', f'Excluded file type: {extension}'
        if extension not in self.include_ext:
            return 'skip', f'Unsupported file type: {extension}'
        return 'analyze', 'Supported file type within retention period'

    def get_file_counts(self, directory_path: Union[str, Path]) -> Dict[str, int]:
        """Get counts of files by category without yielding individual files."""
        counts = {'destroy': 0, 'analyze': 0, 'skip': 0, 'total': 0}
        for file_info in self.scan_directory(directory_path):
            counts[file_info.category] += 1
            counts['total'] += 1
        return counts

def extract_file_content(f: Path, max_chars: int = 4000) -> str:
    """
    Extract text content from a file, using OCR/parsers for binary formats.
    Returns up to max_chars of cleaned text.
    """
    suffix = f.suffix.lower()
    try:
        if suffix == '.txt':
            for encoding in ('utf-8', 'latin-1'):
                try:
                    with f.open('r', encoding=encoding, errors='ignore') as fp:
                        content = fp.read(max_chars)
                        return _clean_text(content)
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    return f"[Error reading file: {str(e)}]"
            return f"[Unreadable file: {suffix}]"
        elif suffix == '.pdf':
            text = ""
            if pdfplumber:
                try:
                    with pdfplumber.open(str(f)) as pdf:
                        for page in pdf.pages:
                            text += page.extract_text() or ""
                            if len(text) >= max_chars:
                                break
                except Exception:
                    text = ""
            if not text and PyPDF2:
                try:
                    with open(f, 'rb') as fp:
                        reader = PyPDF2.PdfReader(fp)
                        for page in reader.pages:
                            text += page.extract_text() or ""
                            if len(text) >= max_chars:
                                break
                except Exception:
                    text = ""
            if not text and pytesseract and Image:
                try:
                    import pdf2image
                    images = pdf2image.convert_from_path(str(f))
                    for img in images:
                        text += pytesseract.image_to_string(img)
                        if len(text) >= max_chars:
                            break
                except Exception:
                    text = ""
            return _clean_text(text)[:max_chars] if text else "[Could not extract text from PDF]"
        elif suffix == '.docx':
            if Document:
                try:
                    doc = Document(str(f))
                    text = "\n".join(p.text for p in doc.paragraphs)
                    return _clean_text(text)[:max_chars]
                except Exception as e:
                    return f"[Error reading DOCX: {str(e)}]"
            else:
                return "[python-docx not installed]"
        elif suffix == '.doc':
            try:
                result = subprocess.run(
                    ["antiword", str(f)], stdout=subprocess.PIPE, check=True
                )
                text = result.stdout.decode("utf-8", errors="ignore")
                return _clean_text(text)[:max_chars]
            except FileNotFoundError:
                return "[antiword not installed for .doc]"
            except subprocess.CalledProcessError as exc:
                return f"[Error reading DOC: {exc}]"
        elif suffix == '.pptx':
            if Presentation:
                try:
                    prs = Presentation(str(f))
                    text: List[str] = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                                if sum(len(t) for t in text) >= max_chars:
                                    break
                        if sum(len(t) for t in text) >= max_chars:
                            break
                    return _clean_text("\n".join(text))[:max_chars]
                except Exception as exc:
                    return f"[Error reading PPTX: {exc}]"
            else:
                return "[python-pptx not installed]"
        elif suffix == '.xlsx':
            if openpyxl:
                try:
                    wb = openpyxl.load_workbook(str(f), read_only=True, data_only=True)
                    text: List[str] = []
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(values_only=True):
                            for cell in row:
                                if cell is not None:
                                    text.append(str(cell))
                                    if sum(len(t) for t in text) >= max_chars:
                                        break
                            if sum(len(t) for t in text) >= max_chars:
                                break
                        if sum(len(t) for t in text) >= max_chars:
                            break
                    return _clean_text(" ".join(text))[:max_chars]
                except Exception as e:
                    return f"[Error reading XLSX: {str(e)}]"
            else:
                return "[openpyxl not installed]"
        else:
            return f"[Unsupported file type: {suffix}]"
    except Exception as e:
        return f"[Error extracting content: {str(e)}]"

def _clean_text(text: str) -> str:
    """Collapse whitespace, strip control chars, etc."""
    import re
    text = re.sub(r'[\r\n]+', '\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()

def main():
    """Test/CLI entry point for the file scanner."""
    import argparse
    parser = argparse.ArgumentParser(description="Scan directory and print file stats.")
    parser.add_argument("directory", help="Directory to scan")
    args = parser.parse_args()
    scanner = FileScanner()
    counts = scanner.get_file_counts(args.directory)
    print("File counts:", counts)
    print("Sample 'analyze' files:")
    for fi in scanner.scan_directory(args.directory):
        if fi.category == "analyze":
            print(f"- {fi.path} ({fi.size_bytes} bytes) | {fi.reason}")

if __name__ == "__main__":
    main()
